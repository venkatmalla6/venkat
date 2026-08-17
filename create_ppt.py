from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MediVision AI"
    subtitle.text = "An Intelligent Multimodal Medical Diagnosis and Clinical Decision Support Platform\nB.Tech Final Year Project"

    # Slide 2: Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Introduction & Objective"
    tf = body_shape.text_frame
    tf.text = "Context: Medical diagnosis is complex and requires doctors to synthesize diverse data (X-rays, lab reports, symptoms, history)."
    p = tf.add_paragraph()
    p.text = "The Challenge: Manual integration is time-consuming and prone to human error, potentially delaying treatment."
    p = tf.add_paragraph()
    p.text = "Project Objective: To develop 'MediVision AI', a multimodal AI platform that streamlines the diagnostic workflow, enhances decision accuracy, and improves patient outcomes."

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Problem Statement"
    tf = body_shape.text_frame
    tf.text = "Information Overload: Doctors face massive amounts of disjointed patient data:"
    p = tf.add_paragraph()
    p.text = " - Medical Images (X-rays, MRIs, CT scans)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = " - Laboratory Reports (Blood tests in PDF format)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = " - Unstructured Documents (Prescriptions, clinical notes)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = " - Patient Symptoms (Verbal or text descriptions)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = " - Historical Records (Past diseases, allergies)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Impact: High risk of overlooking critical details, leading to diagnostic errors or delays."

    # Slide 4: Proposed Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Proposed Solution: MediVision AI"
    tf = body_shape.text_frame
    tf.text = "A unified, AI-powered platform capable of processing multiple data types simultaneously:"
    p = tf.add_paragraph()
    p.text = "Vision: Interprets medical imaging using deep learning."
    p = tf.add_paragraph()
    p.text = "NLP/OCR: Extracts structured data from text and PDF reports."
    p = tf.add_paragraph()
    p.text = "RAG (Retrieval-Augmented Generation): Answers complex medical queries using verified medical knowledge bases."
    p = tf.add_paragraph()
    p.text = "Voice: Multilingual assistant for natural patient interactions."

    # Slide 5: System Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Overall System Architecture"
    tf = body_shape.text_frame
    tf.text = "Data Processing Layer: Handles OCR, Speech-to-Text, and raw data ingestion."
    p = tf.add_paragraph()
    p.text = "Medical AI Layer: Runs Vision Transformers for images and predictive models for disease detection."
    p = tf.add_paragraph()
    p.text = "Knowledge Retrieval (RAG): Connects to medical guidelines and research papers to ground the AI's conclusions."
    p = tf.add_paragraph()
    p.text = "Output Layer: Generates Clinical Summaries, Diagnosis Recommendations, and populates the Doctor Dashboard."

    # Slide 6: Core AI Modules
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Core AI Modules"
    tf = body_shape.text_frame
    tf.text = "Medical Image Analysis:"
    p = tf.add_paragraph()
    p.text = "Analyzes X-rays, MRIs, CT scans."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Detects Pneumonia, Tumors, Fractures, etc."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Blood Reports & OCR:"
    p = tf.add_paragraph()
    p.text = "Extracts key values from PDF reports using Tesseract/PaddleOCR."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Medical RAG System:"
    p = tf.add_paragraph()
    p.text = "Prevents 'AI hallucinations' by retrieving facts from WHO/NIH guidelines."
    p.level = 1

    # Slide 7: Tech Stack & Future
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Tech Stack & Future Scope"
    tf = body_shape.text_frame
    tf.text = "Technology Stack:"
    p = tf.add_paragraph()
    p.text = "Frontend: React, TypeScript, Tailwind CSS"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend/AI: FastAPI, Python, OpenCV, Hugging Face, LangChain"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Databases/Cloud: PostgreSQL, Redis, FAISS, AWS"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Enhancements:"
    p = tf.add_paragraph()
    p.text = "Real-time ECG analysis and Wearable device integration."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Full Electronic Health Record (EHR) integration."
    p.level = 1

    # Slide 8: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Conclusion"
    tf = body_shape.text_frame
    tf.text = "MediVision AI bridges the gap between disparate medical data and actionable clinical insights."
    p = tf.add_paragraph()
    p.text = "It is designed not to replace doctors, but to act as a powerful Clinical Decision Support System."
    p = tf.add_paragraph()
    p.text = "Reduces cognitive load and improves diagnostic accuracy."

    # Save
    prs.save('MediVision_AI_Presentation.pptx')
    print("Presentation saved as MediVision_AI_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
